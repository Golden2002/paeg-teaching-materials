# -*- coding: utf-8 -*-
"""PPT 真实渲染落盘（python-pptx）——render_outline_to_pptx。

移植自主项目 `pptx_mcp_server.generate_ppt`（v0.60 ⭐），剥离宿主的
jieba / requests / 品牌 Logo 依赖，保留核心能力：
- markdown 清理（** / ## / ` / [t](url) / 列表符号，保留内容）
- 品牌配色 + 封面页 + 内容页（统一头部标题条 + 页脚 + 页码）
- 长文本自适应字号（18 → 16 → 14 → 12pt）+ 超长拆页

依赖：python-pptx（可选 extras "pptx"）。未安装时抛 RuntimeError，
由上层生成器捕获并优雅降级（返回提示而非崩溃）。
"""

from __future__ import annotations

import os
import re
from typing import List

from ._paths import default_out_dir, safe_stem


def clean_md(text) -> str:
    """清除 markdown 符号，保留内容。"""
    if not text:
        return text or ""
    t = str(text)
    t = re.sub(r"\*\*(.+?)\*\*", r"\1", t, flags=re.S)
    t = re.sub(r"\*(.+?)\*", r"\1", t, flags=re.S)
    t = re.sub(r"^#{1,6}\s*", "", t, flags=re.M)
    t = re.sub(r"`([^`]*)`", r"\1", t)
    t = re.sub(r"\[([^\]]*)\]\([^)]*\)", r"\1", t)
    t = re.sub(r"^\s*[-*]\s+", "\u2022 ", t, flags=re.M)
    return t.strip()


def _parse_outline(text) -> List[dict]:
    """把 LLM/文本大纲解析为 [{title, points:[...], notes}]。

    兼容 list[dict{slide,title,points}] 输入与 markdown 文本输入。
    """
    if isinstance(text, list):
        slides = []
        for it in text:
            if not isinstance(it, dict):
                continue
            title = str(it.get("title") or it.get("slide") or "")
            pts = it.get("points") or []
            if isinstance(pts, str):
                pts = [ln.strip() for ln in pts.splitlines() if ln.strip()]
            if not title and not pts:
                continue
            slides.append({
                "title": clean_md(title),
                "points": [clean_md(str(p)) for p in pts],
                "notes": str(it.get("notes") or ""),
            })
        if slides:
            return slides
        return [{"title": "演示文稿", "points": ["（空大纲，请补充内容）"], "notes": ""}]

    slides = []
    lines = [l.strip() for l in (text or "").splitlines() if l.strip()]
    cur = None
    for ln in lines:
        m = re.match(r"^(#{1,4})\s+(.+)$", ln)
        if m:
            if cur:
                slides.append(cur)
            cur = {"title": clean_md(m.group(2).strip()), "points": [], "notes": ""}
            continue
        m = re.match(r"^(\d+)[.、)]\s*(.+)$", ln)
        if m:
            if cur:
                slides.append(cur)
            cur = {"title": clean_md(m.group(2).strip()), "points": [], "notes": ""}
            continue
        m = re.match(r"^[-*•]\s*(.+)$", ln)
        if m:
            if cur:
                cur["points"].append(clean_md(m.group(1).strip()))
            else:
                cur = {"title": "要点", "points": [clean_md(m.group(1).strip())], "notes": ""}
            continue
        if cur:
            if ln.startswith("备注") or ln.startswith("note"):
                cur["notes"] = ln.split(":", 1)[-1].strip()
            else:
                cur["points"].append(clean_md(ln))
    if cur:
        slides.append(cur)
    return slides or [{"title": "演示文稿", "points": ["（空大纲，请补充内容）"], "notes": ""}]


# 配色方案（与主项目路演 PPT 一致的品牌色，可按 style 覆盖）
_STYLES = {
    "paeg_standard": {"primary": (0x0F, 0x2A, 0x52), "accent": (0xE6, 0xA5, 0x28),
                      "light": (0xF5, 0xF2, 0xEC), "dark": (0x0F, 0x2A, 0x52)},
    "presentation_zen": {"primary": (0x2C, 0x5F, 0x2D), "accent": (0x97, 0xBC, 0x62),
                         "light": (0xF5, 0xF5, 0xF5), "dark": (0x21, 0x21, 0x21)},
    "dark_premium": {"primary": (0x1E, 0x27, 0x61), "accent": (0xCA, 0xDC, 0xFC),
                     "light": (0x11, 0x14, 0x20), "dark": (0xFF, 0xFF, 0xFF)},
}


def render_outline_to_pptx(outline, topic: str, out_dir=None,
                           style: str = "paeg_standard") -> str:
    """渲染 PPT 大纲 → 真实 .pptx 文件，返回绝对路径。

    Args:
        outline: LLM 生成的大纲（markdown 文本或 list[dict]）。
        topic: 演示主题（用于封面标题与文件名）。
        out_dir: 输出目录（默认 <plugin_root>/render_output/ppt/）。
        style: 'paeg_standard' / 'presentation_zen' / 'dark_premium'。

    Returns:
        str: .pptx 绝对路径。

    Raises:
        RuntimeError: python-pptx 未安装或渲染失败（上层优雅降级）。
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
    except Exception as e:
        raise RuntimeError(f"python-pptx 未安装（pip install python-pptx）: {e}")

    colors = _STYLES.get(style, _STYLES["paeg_standard"])
    C_PRIMARY = RGBColor(*colors["primary"])
    C_ACCENT = RGBColor(*colors["accent"])
    C_LIGHT = RGBColor(*colors["light"])
    C_DARK = RGBColor(*colors["dark"])
    C_GRAY = RGBColor(0x55, 0x5F, 0x6B)
    C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)

    def _set_text(shape, text, size=18, bold=False, color=C_DARK,
                  align=PP_ALIGN.LEFT):
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = str(text)
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = "Microsoft YaHei"
        p.alignment = align
        return shape

    def _add_header(slide, title, page_no, prs):
        bar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(1.0))
        bar.fill.solid()
        bar.fill.fore_color.rgb = C_PRIMARY
        bar.line.fill.background()
        tf = bar.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.5)
        p = tf.paragraphs[0]
        p.text = str(title)
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = C_WHITE
        p.font.name = "Microsoft YaHei"
        p.alignment = PP_ALIGN.LEFT
        _set_text(slide.shapes.add_textbox(Inches(0.5), Inches(7.1), Inches(6), Inches(0.3)),
                  "PAEG · 教学物料", size=8, color=C_GRAY)
        _set_text(slide.shapes.add_textbox(Inches(12.2), Inches(7.1), Inches(1.0), Inches(0.3)),
                  str(page_no), size=8, color=C_GRAY, align=PP_ALIGN.RIGHT)

    def _add_bullets_adaptive(slide, points, left, top, width, height):
        full_text = "\n".join(str(p) for p in points)
        total_chars = len(full_text)
        size = 18
        if total_chars > 400:
            size = 16
        if total_chars > 700:
            size = 14
        if total_chars > 1100:
            size = 12
        chunks = [points[i:i + 6] for i in range(0, len(points), 6)]
        for ci, chunk in enumerate(chunks):
            box = slide.shapes.add_textbox(Inches(left), Inches(top),
                                           Inches(width), Inches(height))
            tf = box.text_frame
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.NONE
            for i, pt in enumerate(chunk):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                pt = clean_md(str(pt))
                p.text = ("\u2022 " if not pt.startswith("\u2022") else "") + pt
                p.font.size = Pt(size)
                p.font.color.rgb = C_DARK
                p.font.name = "Microsoft YaHei"
                p.space_after = Pt(8)
        return box

    try:
        out_dir = out_dir or default_out_dir("ppt")
        os.makedirs(out_dir, exist_ok=True)
        slides_data = _parse_outline(outline or topic)

        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        blank = prs.slide_layouts[6]

        # 封面页
        s = prs.slides.add_slide(blank)
        bg = s.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = C_PRIMARY
        bg.line.fill.background()
        _set_text(s.shapes.add_textbox(Inches(1.2), Inches(2.0), Inches(11), Inches(1.6)),
                  topic, size=40, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        _set_text(s.shapes.add_textbox(Inches(1.2), Inches(3.8), Inches(11), Inches(1.0)),
                  "PAEG · 教学演示文稿", size=18, color=C_LIGHT, align=PP_ALIGN.CENTER)

        # 内容页
        for i, sd in enumerate(slides_data[:20]):
            s = prs.slides.add_slide(blank)
            _add_header(s, sd["title"], i + 2, prs)
            pts = sd["points"] or ["（本页要点）"]
            _add_bullets_adaptive(s, pts, 0.8, 1.4, 11.7, 5.4)
            if sd.get("notes"):
                try:
                    s.notes_slide.notes_text_frame.text = sd["notes"]
                except Exception:
                    pass

        fname = safe_stem(topic, maxlen=40) + ".pptx"
        path = os.path.join(out_dir, fname)
        prs.save(path)
        return os.path.abspath(path)
    except RuntimeError:
        raise
    except Exception as e:
        raise RuntimeError(f"PPT 渲染失败: {str(e)[:200]}")
